import os
import pypdf
import docx
import pptx
import traceback

class DocumentIndexer:
    SUPPORTED_EXTENSIONS = {
        '.pdf': 'PDF',
        '.docx': 'Word (DOCX)',
        '.pptx': 'PowerPoint (PPTX)',
        '.txt': 'Text (TXT)',
        '.md': 'Markdown (MD)',
        '.csv': 'CSV',
        '.log': 'Log (LOG)',
        '.json': 'JSON',
        '.py': 'Python (PY)',
        '.js': 'JavaScript (JS)',
        '.jsx': 'React (JSX)',
        '.ts': 'TypeScript (TS)',
        '.tsx': 'TypeScript (TSX)',
        '.html': 'HTML Document',
        '.htm': 'HTML Document',
        '.css': 'CSS Stylesheet',
        '.java': 'Java Source',
        '.kt': 'Kotlin Source',
        '.c': 'C Source',
        '.cpp': 'C++ Source',
        '.h': 'C/C++ Header',
        '.cs': 'C# Source',
        '.rs': 'Rust Source',
        '.go': 'Go Source',
        '.php': 'PHP Source',
        '.rb': 'Ruby Source',
        '.sql': 'SQL Script',
        '.xml': 'XML Document',
        '.yaml': 'YAML Config',
        '.yml': 'YAML Config',
        '.sh': 'Shell Script',
        '.bat': 'Batch Script',
        '.ini': 'Config (INI)',
        '.env': 'Environment Config',
        '.conf': 'Configuration',
        '.toml': 'TOML Config',
        '.rst': 'reStructuredText',
        '.rtf': 'Rich Text (RTF)'
    }
    
    IMAGE_EXTENSIONS = {
        '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp', '.ico', '.svg'
    }

    def __init__(self, db_manager):
        self.db = db_manager
        self._ocr_reader = None

    def get_ocr_reader(self):
        if self._ocr_reader is None:
            try:
                import easyocr
                try:
                    self._ocr_reader = easyocr.Reader(['en'], gpu=True)
                except Exception:
                    self._ocr_reader = easyocr.Reader(['en'], gpu=False)
            except Exception as ocr_init_err:
                print(f"EasyOCR initialization bypassed: {ocr_init_err}")
                self._ocr_reader = False
        return self._ocr_reader if self._ocr_reader is not False else None

    def parse_image(self, file_path):
        try:
            reader = self.get_ocr_reader()
            if not reader:
                return ""
            results = reader.readtext(file_path)
            text_lines = [res[1] for res in results if res[1]]
            return "\n".join(text_lines)
        except Exception as e:
            print(f"Error running OCR on image {file_path}: {e}")
            return ""

    def parse_pdf(self, file_path):
        text_content = []
        try:
            with open(file_path, 'rb') as f:
                reader = pypdf.PdfReader(f)
                for page_num in range(len(reader.pages)):
                    page = reader.pages[page_num]
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
            return "\n".join(text_content)
        except Exception as e:
            print(f"Error parsing PDF {file_path}: {e}")
            return ""

    def parse_docx(self, file_path):
        try:
            doc = docx.Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text]
            # Also extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells if cell.text]
                    if row_text:
                        paragraphs.append(" | ".join(row_text))
            return "\n".join(paragraphs)
        except Exception as e:
            print(f"Error parsing DOCX {file_path}: {e}")
            return ""

    def parse_pptx(self, file_path):
        try:
            prs = pptx.Presentation(file_path)
            slide_texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text)
            return "\n".join(slide_texts)
        except Exception as e:
            print(f"Error parsing PPTX {file_path}: {e}")
            return ""

    def parse_text(self, file_path):
        # Try UTF-8 first, fallback to latin-1 / cp1252 for encoding tolerance
        for encoding in ['utf-8', 'latin-1', 'cp1252', 'utf-16']:
            try:
                with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
            except Exception as e:
                print(f"Error reading text file {file_path} with {encoding}: {e}")
                break
        return ""

    def scan_directory(self, dir_path, user_id='', yield_progress=None):
        """
        Scans a directory tree, parses files, and indexes them into database.
        user_id tags all documents to the requesting user.
        yield_progress is a callback function for streaming progress logs.
        """
        if not os.path.exists(dir_path):
            raise FileNotFoundError(f"Directory {dir_path} does not exist.")
            
        indexed_count = 0
        skipped_count = 0
        error_count = 0
        
        def walk_error(err):
            print(f"Skipping inaccessible system folder: {err.filename} ({err.strerror})")

        all_files = []
        ignore_dirs = {
            'node_modules', 'venv', '__pycache__', 'appdata', 'application data', 
            'cookies', 'local settings', 'sendto', 'start menu', 'templates', 'history',
            'system volume information', '$recycle.bin', 'msocache', 'recycler', 'documents and settings'
        }
        
        for root, dirs, files in os.walk(dir_path, onerror=walk_error):
            # Ignore hidden folders and junctions case-insensitively
            dirs[:] = [d for d in dirs if not d.startswith('.') and d.lower() not in ignore_dirs]
            
            for file in files:
                if file.startswith('~$') or file.startswith('.'):
                    continue
                try:
                    filepath = os.path.join(root, file)
                    ext = os.path.splitext(file)[1].lower()
                    all_files.append((filepath, file, ext))
                except Exception as file_err:
                    print(f"Error processing file path details: {file_err}")
        
        total_files = len(all_files)
        if yield_progress:
            yield_progress(f"Found total of {total_files} files in directory tree.")
            
        for index, (filepath, filename, ext) in enumerate(all_files):
            if ext in self.IMAGE_EXTENSIONS:
                filetype = 'Image OCR'
                try:
                    filesize = os.path.getsize(filepath)
                except Exception:
                    filesize = 0
                if yield_progress:
                    yield_progress(f"[{index+1}/{total_files}] Running OCR on Image: {filename}...")
                content = self.parse_image(filepath)
            elif ext in self.SUPPORTED_EXTENSIONS:
                filetype = self.SUPPORTED_EXTENSIONS[ext]
                try:
                    filesize = os.path.getsize(filepath)
                except Exception:
                    filesize = 0
                if yield_progress:
                    yield_progress(f"[{index+1}/{total_files}] Parsing {filetype}: {filename}...")
                
                content = ""
                if ext == '.pdf':
                    content = self.parse_pdf(filepath)
                elif ext == '.docx':
                    content = self.parse_docx(filepath)
                elif ext == '.pptx':
                    content = self.parse_pptx(filepath)
                else:
                    content = self.parse_text(filepath)
            else:
                skipped_count += 1
                continue
                
            # Check if parsed content is meaningful
            if content.strip():
                try:
                    self.db.upsert_document(
                        user_id=user_id,
                        filepath=filepath,
                        filename=filename,
                        filetype=filetype,
                        filesize=filesize,
                        content_text=content
                    )
                    indexed_count += 1
                except Exception as db_err:
                    error_count += 1
                    if yield_progress:
                        yield_progress(f"Error indexing {filename} in DB: {db_err}")
            else:
                error_count += 1
                if yield_progress:
                    yield_progress(f"Warning: Extracted empty text or failed to read: {filename}")
                    
        if yield_progress:
            yield_progress(f"Scanning completed: {indexed_count} indexed, {skipped_count} skipped/ignored, {error_count} failed.")
            
        return {
            "indexed": indexed_count,
            "skipped": skipped_count,
            "failed": error_count
        }
