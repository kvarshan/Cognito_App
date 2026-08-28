import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

def create_folders():
    folders = [
        "demo_data/Gen AI",
        "demo_data/Business & Finance",
        "demo_data/Software Development"
    ]
    for f in folders:
        os.makedirs(f, exist_ok=True)
    print("Folders created successfully.")

def create_docx_file():
    doc_path = "demo_data/Gen AI/llm_fundamentals.docx"
    doc = Document()
    doc.add_heading("Large Language Models and Transformer Architecture", 0)
    
    doc.add_heading("1. The Transformer Revolution", level=1)
    doc.add_paragraph(
        "Introduced in 2017 in the seminal paper 'Attention Is All You Need', the Transformer architecture "
        "has become the backbone of modern Generative AI. Unlike previous sequence-to-sequence models like "
        "Recurrent Neural Networks (RNNs) and Long Short-Term Memory (LSTM) networks, Transformers process "
        "entire sequences in parallel, dramatically improving training efficiency and scalability."
    )
    
    doc.add_heading("2. Core Mechanisms: Self-Attention", level=2)
    doc.add_paragraph(
        "The key innovation of the Transformer is the self-attention mechanism. It allows the model to "
        "weigh the importance of different words in a sentence, regardless of their distance from one another. "
        "For example, in the sentence 'The bank of the river had a bank that was closed,' the model can "
        "use context to disambiguate the two meanings of the word 'bank'."
    )
    
    doc.add_heading("3. Local Model Training and Fine-Tuning", level=2)
    doc.add_paragraph(
        "Training LLMs locally requires massive computational power, usually leveraging GPUs. "
        "Techniques like Low-Rank Adaptation (LoRA) and Parameter-Efficient Fine-Tuning (PEFT) "
        "have democratized local model training. This allows developers to fine-tune pre-trained models "
        "like Google Gemini, Llama, and Mistral on proprietary data while maintaining complete offline data privacy."
    )
    
    doc.save(doc_path)
    print(f"Created Word Document at {doc_path}")

def create_pptx_file():
    ppt_path = "demo_data/Business & Finance/q1_financials.pptx"
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Q1 Financial & Operations Review"
    subtitle.text = "Cognito Corp - High Privacy Local Intelligence"
    
    # Slide 2: Financial Stats
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Financial Growth Highlights"
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Key Metrics for Q1 2026:"
    
    p = tf.add_paragraph()
    p.text = "- Revenue grew by 24% year-over-year, reaching $4.8M."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "- Gross margins sustained at 78% due to optimized offline infrastructure costs."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "- Local AI processing reduced cloud expenses by 40%."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "- Net cash flow from operating activities was positive at $1.2M."
    p.level = 1
    
    prs.save(ppt_path)
    print(f"Created PowerPoint Presentation at {ppt_path}")

def create_text_and_md_files():
    # Prompt Engineering
    pe_content = """# Guide to Advanced Prompt Engineering

Prompt engineering is the practice of structuring text inputs to Generative AI models to achieve the most accurate and high-quality outputs. Since local offline models have fixed knowledge, prompt quality is critical for performance.

## 1. Core Prompting Techniques
- **Zero-Shot Prompting**: Asking the model a question directly without providing any examples. E.g., "Translate 'Hello' to French."
- **Few-Shot Prompting**: Providing a few high-quality input-output demonstrations before asking the target question. This helps the model align with the desired format and style.
- **Chain of Thought (CoT)**: Instructing the model to break down its reasoning step-by-step before answering. This is highly effective for mathematical, logical, and symbolic reasoning tasks.

## 2. Best Practices for System Prompts
- Define the persona and role explicitly (e.g., "You are an expert AI software architect").
- Specify constraints clearly (e.g., "Use only offline resources," "Keep the response under 3 sentences").
- Use structured formats (Markdown tables, JSON outputs, bullet points).
"""
    with open("demo_data/Gen AI/prompt_engineering.md", "w", encoding="utf-8") as f:
        f.write(pe_content)
        
    # Agentic Workflows
    agents_content = """# Agentic Workflows and Multi-Agent Systems

Generative AI is shifting from static, conversational chatbots to proactive, agentic workflows. An "AI Agent" is a model equipped with tools (calculators, web browsers, local files, search engines) that operates in a loop: Plan -> Act -> Observe.

## Key Properties of Agents
1. **Autonomy**: Ability to decide its own plan of action based on a high-level goal.
2. **Tool Use**: Triggering custom functions or commands to gather external information.
3. **Memory**: Retaining historical execution logs (short-term) and semantic databases (long-term).
4. **Execution Loops**: Iteratively reviewing outcomes and self-correcting mistakes.

By training local search indexes like Cognito, local agents can rapidly browse files offline without exposing sensitive documents to third-party APIs.
"""
    with open("demo_data/Gen AI/agentic_workflows.md", "w", encoding="utf-8") as f:
        f.write(agents_content)
        
    # Market Forecast
    mf_content = """Market Trends and Strategic Forecast 2026

Executive Summary:
The global macroeconomic environment in 2026 is heavily shaped by data privacy regulations, high cloud computing energy costs, and the rise of edge intelligence. Large-scale cloud services face increasing pushback due to privacy leaks and corporate espionage concerns.

Strategic Decisions:
1. Shift to Edge and Local AI: Our organization must invest in local document parsing and indexing. All training, classification, and searching should occur on-premise.
2. Reduction of SaaS dependencies: Moving critical data (Dox, PDF, PPT, TXT) to local storage.
3. Adoption of Hybrid Search: Implementing localized keyword and cluster-based search systems to query offline storage securely.
"""
    with open("demo_data/Business & Finance/market_forecast.txt", "w", encoding="utf-8") as f:
        f.write(mf_content)

    # Python Quickstart
    py_content = """# Python Quickstart Guide for Machine Learning

Python is the absolute standard for data science and AI development. Below is a quick cheatsheet of essential concepts used in building local indexes and ML models.

## 1. List and Dict Comprehensions
Comprehensions are concise ways to generate lists and dictionaries.
```python
# List comprehension
squared_numbers = [x**2 for x in range(10)]

# Dictionary comprehension
document_lengths = {doc_name: len(text) for doc_name, text in docs.items()}
```

## 2. Text Preprocessing and Tokenization
Before training vectorizers, text must be normalized.
```python
def preprocess(text):
    # Lowercase and split into words
    words = text.lower().split()
    # Filter non-alphabetic tokens
    return [w for w in words if w.isalpha()]
```

## 3. Object-Oriented Database Connectors
Wrapping SQLite connections in clear classes.
```python
import sqlite3

class LocalDB:
    def __init__(self, path):
        self.conn = sqlite3.connect(path)
        
    def query(self, sql, params=()):
        with self.conn:
            return self.conn.execute(sql, params).fetchall()
```
"""
    with open("demo_data/Software Development/python_quickstart.md", "w", encoding="utf-8") as f:
        f.write(py_content)
        
    # API Design
    api_content = """RESTful API Design & Best Practices

REST (Representational State Transfer) is an architectural style for designing networked applications. It relies on stateless, client-server communication using HTTP protocols.

Key HTTP Verbs:
- GET: Retrieve a resource (e.g., GET /api/documents)
- POST: Create a new resource or perform an action (e.g., POST /api/train)
- PUT: Update an existing resource entirely
- PATCH: Update parts of an existing resource
- DELETE: Remove a resource

Common HTTP Status Codes:
- 200 OK: Request succeeded
- 201 Created: Resource successfully created
- 400 Bad Request: Invalid input from client
- 401 Unauthorized: Authentication required
- 404 Not Found: Resource does not exist
- 500 Internal Server Error: Backend crash

In Cognito, we use FastAPI to expose lightweight endpoints that serve local document data to our Single Page Application.
"""
    with open("demo_data/Software Development/api_design.txt", "w", encoding="utf-8") as f:
        f.write(api_content)
        
    print("Text and Markdown files created successfully.")

if __name__ == "__main__":
    create_folders()
    create_docx_file()
    create_pptx_file()
    create_text_and_md_files()
    print("Demo dataset built successfully!")
